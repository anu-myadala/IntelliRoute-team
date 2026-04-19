#!/usr/bin/env python3
"""Generate IntelliRoute presentation slides (10 slides, .pptx)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Colour palette ────────────────────────────────────────────────
BG       = RGBColor(0x14, 0x14, 0x1E)   # dark navy
ACCENT   = RGBColor(0x64, 0x9C, 0xFF)   # bright blue
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0xAA, 0xAA, 0xBB)
GREEN    = RGBColor(0x4E, 0xC9, 0xB0)
ORANGE   = RGBColor(0xFF, 0xA5, 0x4C)
RED      = RGBColor(0xFF, 0x6B, 0x6B)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def _set_bg(slide, colour=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colour


def _add_text(slide, left, top, width, height, text, *,
              size=18, bold=False, colour=WHITE, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.paragraphs[0].alignment = align
    except Exception:
        pass
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = colour
    return txBox


def _add_bullets(slide, left, top, width, height, items, *,
                 size=16, colour=WHITE, bullet_colour=ACCENT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(6)
        p.space_after = Pt(2)

        # bullet character
        bullet_run = p.add_run()
        bullet_run.text = "\u25B8 "
        bullet_run.font.size = Pt(size)
        bullet_run.font.color.rgb = bullet_colour
        bullet_run.font.bold = True

        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = colour
    return txBox


def _title_bar(slide, title, subtitle=None):
    _add_text(slide, 0.6, 0.3, 12, 0.7, title,
              size=36, bold=True, colour=ACCENT)
    if subtitle:
        _add_text(slide, 0.6, 1.0, 12, 0.5, subtitle,
                  size=18, colour=GRAY)


def _add_table(slide, left, top, width, height, headers, rows, *,
               hdr_colour=ACCENT, cell_colour=WHITE, bg_colour=RGBColor(0x1E, 0x1E, 0x2E)):
    tbl_shape = slide.shapes.add_table(
        len(rows) + 1, len(headers),
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tbl = tbl_shape.table

    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = hdr
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(14)
                r.font.bold = True
                r.font.color.rgb = hdr_colour
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2A)

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(13)
                    r.font.color.rgb = cell_colour
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_colour if ri % 2 == 0 else RGBColor(0x24, 0x24, 0x34)
    return tbl_shape


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
_set_bg(s)
_add_text(s, 1.5, 1.5, 10, 1.2,
          "IntelliRoute",
          size=60, bold=True, colour=ACCENT, align=PP_ALIGN.CENTER)
_add_text(s, 1.5, 2.8, 10, 0.8,
          "A Distributed Control Plane for Multi-LLM Orchestration",
          size=24, colour=WHITE, align=PP_ALIGN.CENTER)
_add_text(s, 1.5, 4.2, 10, 0.5,
          "CMPE 273 \u2014 Enterprise Distributed Systems  |  Spring 2026",
          size=18, colour=GRAY, align=PP_ALIGN.CENTER)
_add_text(s, 1.5, 5.2, 10, 0.5,
          "Anukrithi Myadala  \u00B7  Larry Nguyen  \u00B7  James Pham  \u00B7  Surbhi Singh",
          size=18, colour=GRAY, align=PP_ALIGN.CENTER)
_add_text(s, 1.5, 5.8, 10, 0.5,
          "San Jos\u00E9 State University",
          size=16, colour=GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(s)
_title_bar(s, "The Problem", "Why enterprises need a multi-LLM control plane")
_add_bullets(s, 0.8, 1.7, 5.5, 5.0, [
    "Organisations use multiple LLM providers (OpenAI, Anthropic, open-source)",
    "Each has different latency, cost, rate limits, and failure modes",
    "Every application team reinvents: retries, fallback, cost tracking",
    "No unified place to make routing decisions across heterogeneous workloads",
    "Costs escalate without per-tenant visibility or governance",
], size=17)
_add_text(s, 7.0, 1.7, 5.5, 5.0,
          "Without IntelliRoute:\n\n"
          "  App A \u2192 OpenAI    (own retry logic)\n"
          "  App B \u2192 Anthropic (own rate limiter)\n"
          "  App C \u2192 Both      (hardcoded fallback)\n\n"
          "With IntelliRoute:\n\n"
          "  App A \u2510\n"
          "  App B \u2524\u2500\u2192 IntelliRoute \u2500\u2192 Best Provider\n"
          "  App C \u2518\n\n"
          "  One control plane. Intent-aware. Self-healing.",
          size=16, colour=GREEN)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — Architecture
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(s)
_title_bar(s, "System Architecture", "6 services, all Python/FastAPI, communicating over HTTP")
_add_text(s, 0.8, 1.7, 11.5, 4.5,
          "Client\n"
          "  \u2502\n"
          "  \u25BC\n"
          "Gateway (:8000)  \u2500 API-key auth, tenant rewriting, X-Request-Id\n"
          "  \u2502\n"
          "  \u25BC\n"
          "Router (:8001)   \u2500 Intent classify \u2192 Policy rank \u2192 Fallback loop\n"
          "  \u251C\u2500\u2500sync\u2500\u2500\u25B6 Rate Limiter (:8002)   Token-bucket, leader election (3 replicas)\n"
          "  \u251C\u2500\u2500sync\u2500\u2500\u25B6 Health Monitor (:8004) Circuit breakers, liveness polling\n"
          "  \u2514\u2500\u2500async\u2500\u25B6 Cost Tracker (:8003)   Fire-and-forget events, per-tenant rollups\n"
          "  \u2502\n"
          "  \u25BC\n"
          "Mock LLM Providers\n"
          "  mock-fast (:9001)   30ms, $0.002/1K   \u2500 interactive, code\n"
          "  mock-smart (:9002)  120ms, $0.02/1K   \u2500 reasoning\n"
          "  mock-cheap (:9003)  80ms, $0.0003/1K  \u2500 batch",
          size=16, colour=WHITE)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — Request Lifecycle
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(s)
_title_bar(s, "Request Lifecycle", "What happens when a prompt arrives")
_add_bullets(s, 0.8, 1.7, 11.5, 5.5, [
    'Gateway authenticates X-API-Key \u2192 maps to tenant_id (overwrites body for security)',
    'Gateway generates X-Request-Id (UUID4) for distributed tracing',
    'Router classifies intent: interactive / reasoning / batch / code',
    'Router fetches health snapshot from Health Monitor (circuit breaker states)',
    'Policy engine scores & ranks all providers (4 axes, intent-specific weights)',
    'Fallback loop: check rate limit \u2192 call provider \u2192 if fail, try next',
    'On success: record feedback EMA, publish async cost event, return response',
    'Response includes: provider, model, latency, cost, fallback_used, degraded flags',
], size=16)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — Distributed Systems Concepts
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(s)
_title_bar(s, "Distributed Systems Concepts")
_add_table(s, 0.6, 1.5, 12, 5.5,
    ["Course Topic", "Implementation"],
    [
        ["Service Discovery", "ProviderRegistry \u2014 in-memory analogue of Consul/etcd"],
        ["Sync + Async Comm", "HTTP for routing decisions; fire-and-forget for cost events"],
        ["Leader Election", "Bully algorithm \u2014 3 rate-limiter replicas, highest ID wins"],
        ["Consistency", "Strong for rate limits (single leader); eventual for costs"],
        ["Fault Tolerance", "Circuit breakers (3-state) + automatic fallback chain"],
        ["Backpressure", "Priority queue with load shedding (shed LOW at 80% depth)"],
        ["Multi-Objective Opt.", "Weighted scoring: latency, cost, capability, success rate"],
        ["Security", "API-key auth at gateway; server-side tenant identity rewriting"],
        ["Observability", "Structured JSON logs, /snapshot endpoints, X-Request-Id tracing"],
    ]
)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — Intent-Aware Routing Policy
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(s)
_title_bar(s, "Intent-Aware Routing Policy",
           "Multi-objective scoring with intent-specific weight vectors")

_add_text(s, 0.8, 1.6, 5.5, 0.5, "Intent Classification (Heuristic):",
          size=18, bold=True, colour=GREEN)
_add_bullets(s, 0.8, 2.1, 5.5, 2.5, [
    'Batch: "summarize the following", "extract", "batch"',
    'Code: "```", "def ", "traceback", regex(bug|error)',
    'Reasoning: "explain" + "step by step" + >200 chars',
    'Default: Interactive',
], size=14)

_add_text(s, 0.8, 4.2, 5.5, 0.5, "Score = \u03A3(weight\u2097 \u00D7 score\u2097) \u2212 anomaly_penalty",
          size=16, bold=True, colour=ORANGE)

_add_table(s, 6.8, 1.6, 5.8, 3.0,
    ["Intent", "Latency", "Cost", "Capability", "Success"],
    [
        ["Interactive", "0.55", "0.15", "0.20", "0.10"],
        ["Reasoning",   "0.10", "0.10", "0.50", "0.30"],
        ["Batch",       "0.05", "0.65", "0.15", "0.15"],
        ["Code",        "0.25", "0.10", "0.45", "0.20"],
    ]
)

_add_text(s, 6.8, 5.0, 5.8, 1.5,
          "Result:\n"
          "  \u25B8 Interactive \u2192 mock-fast (latency-optimised)\n"
          "  \u25B8 Reasoning  \u2192 mock-smart (capability-optimised)\n"
          "  \u25B8 Batch      \u2192 mock-cheap (cost-optimised)",
          size=15, colour=WHITE)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — Fault Tolerance
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(s)
_title_bar(s, "Fault Tolerance & Self-Healing")

_add_text(s, 0.8, 1.6, 5.5, 0.5, "Circuit Breaker (3-State Machine):",
          size=18, bold=True, colour=GREEN)
_add_text(s, 0.8, 2.2, 5.8, 2.5,
          "CLOSED \u2500\u2500(3 failures)\u2500\u2500\u25B6 OPEN\n"
          "  \u25B2                              \u2502\n"
          "  \u2502                         (5s timeout)\n"
          "  \u2502                              \u2502\n"
          "  \u2514\u2500\u2500(2 successes)\u2500\u2500 HALF_OPEN \u25C0\u2500\u2518\n"
          "        \u2502\n"
          "  (any failure) \u2500\u2500\u25B6 back to OPEN",
          size=15, colour=WHITE)

_add_text(s, 0.8, 5.0, 5.5, 0.5, "Sliding window: 20 outcomes",
          size=14, colour=GRAY)

_add_text(s, 7.0, 1.6, 5.5, 0.5, "Automatic Fallback Chain:",
          size=18, bold=True, colour=GREEN)
_add_bullets(s, 7.0, 2.2, 5.5, 4.5, [
    "Router ranks all providers by policy score",
    "Iterates through ranked list:",
    "  \u2192 Check rate limit (skip if exhausted)",
    "  \u2192 Check circuit breaker (skip if open)",
    "  \u2192 Call provider (5s timeout)",
    "  \u2192 If fail: report to health monitor, try next",
    "Response flags: fallback_used=true, degraded=true",
    "All providers down? Degraded mode \u2192 structured 503",
], size=15)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 8 — Rate Limiting & Leader Election
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(s)
_title_bar(s, "Distributed Rate Limiting & Leader Election")

_add_text(s, 0.8, 1.6, 5.5, 0.5, "Token Bucket Algorithm:",
          size=18, bold=True, colour=GREEN)
_add_bullets(s, 0.8, 2.2, 5.5, 2.5, [
    "Capacity = 10 tokens (burst size)",
    "Refill rate = 1 token/sec (60 req/min steady-state)",
    "Per (tenant, provider) pair \u2014 independent buckets",
    "Continuous refill: tokens += elapsed \u00D7 rate (capped)",
    "Denied? Returns retry_after_ms to client",
], size=15)

_add_text(s, 0.8, 4.6, 5.5, 0.5, "Why token bucket?",
          size=16, bold=True, colour=ORANGE)
_add_text(s, 0.8, 5.1, 5.5, 1.5,
          "Fixed windows have the boundary problem (2x burst\n"
          "at window edge). Token bucket allows controlled bursts\n"
          "while enforcing a smooth steady-state rate.",
          size=14, colour=GRAY)

_add_text(s, 7.0, 1.6, 5.5, 0.5, "Bully Leader Election:",
          size=18, bold=True, colour=GREEN)
_add_bullets(s, 7.0, 2.2, 5.5, 2.0, [
    "3 replicas: rl-0 (:8002), rl-1 (:8012), rl-2 (:8022)",
    "Highest ID always wins (rl-2 = leader)",
    "Leader sends heartbeats every 1 second",
    "Followers: no heartbeat for 3s \u2192 trigger re-election",
], size=15)

_add_text(s, 7.0, 4.2, 5.5, 0.5, "Consistency Model:",
          size=16, bold=True, colour=ORANGE)
_add_table(s, 7.0, 4.8, 5.5, 2.0,
    ["Data Path", "Consistency"],
    [
        ["Rate-limit tokens", "Strong (single leader)"],
        ["Cost rollups", "Eventual (async fire-and-forget)"],
        ["Circuit breakers", "Per-instance"],
    ]
)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 9 — Testing & Observability
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(s)
_title_bar(s, "Testing & Observability", "42 tests \u2014 35 unit + 7 integration")

_add_table(s, 0.6, 1.6, 6.0, 4.5,
    ["Test Suite", "Tests", "Coverage"],
    [
        ["test_intent.py",          "7", "Intent classifier edge cases"],
        ["test_policy.py",          "7", "Scoring + ranking verification"],
        ["test_token_bucket.py",    "7", "Refill, consume, edge cases"],
        ["test_circuit_breaker.py", "6", "State transitions"],
        ["test_registry.py",       "4", "Register / deregister / lookup"],
        ["test_accounting.py",     "4", "Rollups + budget alerts"],
        ["test_feedback.py",      "~7", "EMA calculations, anomaly"],
        ["test_queue.py",         "~7", "Priority, shedding, caps"],
        ["test_election.py",      "~7", "Election, heartbeats"],
        ["test_integration.py",    "7", "Full 8-process stack over HTTP"],
    ]
)

_add_text(s, 7.2, 1.6, 5.5, 0.5, "Observability:",
          size=18, bold=True, colour=GREEN)
_add_bullets(s, 7.2, 2.2, 5.5, 4.0, [
    "Structured JSON logging (machine-parseable)",
    "X-Request-Id tracing across all services",
    "/snapshot endpoints for circuit breaker state",
    "/feedback for per-provider EMA metrics",
    "/queue/stats for backpressure monitoring",
    "/election/status on each rate-limiter replica",
    "/v1/cost/summary for per-tenant spend",
    "Live dashboard (port 3000) polls all endpoints",
], size=15)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 10 — Live Demo & Future Work
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(s)
_title_bar(s, "Live Demo & Future Work")

_add_text(s, 0.8, 1.6, 5.5, 0.5, "Demo Flow:",
          size=18, bold=True, colour=GREEN)
_add_bullets(s, 0.8, 2.2, 5.5, 4.0, [
    'Start stack: python3 scripts/start_stack.py',
    'Open dashboard: http://localhost:3000',
    'Send "Hi there!" \u2192 routes to mock-fast',
    'Send reasoning prompt \u2192 routes to mock-smart',
    'Send "Summarize..." \u2192 routes to mock-cheap',
    'Force-fail mock-fast \u2192 auto-failover demo',
    'Tighten rate limit \u2192 backpressure demo',
    'Show cost panel + leader election panel',
    'Run pytest \u2014 42 tests pass',
], size=15)

_add_text(s, 7.0, 1.6, 5.5, 0.5, "Production Extensions:",
          size=18, bold=True, colour=GREEN)
_add_bullets(s, 7.0, 2.2, 5.5, 2.0, [
    "Real Raft consensus for rate-limiter replication",
    "Persistent storage (Redis / Postgres) \u2014 single-file swap",
    "Learned intent classifier (fine-tuned small model)",
    "Real LLM provider adapters (OpenAI, Anthropic, vLLM)",
], size=15)

_add_text(s, 7.0, 4.5, 5.5, 0.5, "Key Takeaway:",
          size=18, bold=True, colour=ORANGE)
_add_text(s, 7.0, 5.1, 5.5, 1.5,
          "IntelliRoute treats LLMs as distributed compute\n"
          "resources behind a control plane \u2014 not as isolated\n"
          "APIs. The same patterns (circuit breakers, leader\n"
          "election, backpressure) that power Netflix and\n"
          "Uber apply to the LLM orchestration problem.",
          size=15, colour=WHITE)


# ── Save ──────────────────────────────────────────────────────────
out = "IntelliRoute_Presentation.pptx"
prs.save(out)
print(f"Saved {out}  ({len(prs.slides)} slides)")
